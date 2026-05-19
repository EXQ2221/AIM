from __future__ import annotations

import base64
import io
import json
import os
import re
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

import requests
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pypdf import PdfReader


MAX_DESCRIBE_IMAGES = int(os.getenv("PARSER_MAX_DESCRIBE_IMAGES", "8"))
IMAGE_HEAVY_MIN_IMAGES = int(os.getenv("PARSER_IMAGE_HEAVY_MIN_IMAGES", "4"))
IMAGE_HEAVY_TEXT_CHARS_PER_IMAGE = int(os.getenv("PARSER_IMAGE_HEAVY_TEXT_CHARS_PER_IMAGE", "800"))

VISION_BASE_URL = os.getenv("VISION_BASE_URL", "").strip()
VISION_API_KEY = os.getenv("VISION_API_KEY", "").strip()
VISION_MODEL = os.getenv("VISION_MODEL", "qwen3.6-plus").strip()
VISION_TIMEOUT_SECONDS = int(os.getenv("VISION_TIMEOUT_SECONDS", "90"))


@dataclass
class ParsedDocument:
    title: str
    source_type: str
    content: str
    file_type: str
    image_count: int
    used_vision_description: bool


def parse_document(filename: str, content_type: str, data: bytes, title_override: str | None = None) -> ParsedDocument:
    if not data:
        raise ValueError("file is empty")

    title = (title_override or "").strip() or Path(filename or "").stem.strip() or "Imported Document"
    file_type = detect_file_type(filename, content_type)

    if file_type == "text":
        content = normalize_text(decode_utf8(data))
        return ParsedDocument(title, "TEXT", ensure_not_empty(content), file_type, 0, False)
    if file_type == "markdown":
        content = normalize_text(decode_utf8(data))
        return ParsedDocument(title, "MARKDOWN", ensure_not_empty(content), file_type, 0, False)
    if file_type == "pdf":
        text, images = parse_pdf(data)
        return build_multimodal_document(title, file_type, text, images)
    if file_type == "pptx":
        text, images = parse_pptx(data)
        return build_multimodal_document(title, file_type, text, images)
    if file_type == "docx":
        text = parse_docx(data)
        return ParsedDocument(title, "TEXT", ensure_not_empty(text), file_type, 0, False)
    if file_type == "ppt":
        raise ValueError("legacy .ppt is not supported yet, please convert it to .pptx")
    if file_type == "doc":
        raise ValueError("legacy .doc is not supported yet, please convert it to .docx")
    raise ValueError("unsupported document type, only txt/md/pdf/docx/pptx are supported")


def detect_file_type(filename: str, content_type: str) -> str:
    ext = Path(filename or "").suffix.lower()
    if ext in {".txt"}:
        return "text"
    if ext in {".md", ".markdown"}:
        return "markdown"
    if ext == ".pdf":
        return "pdf"
    if ext == ".docx":
        return "docx"
    if ext == ".doc":
        return "doc"
    if ext == ".pptx":
        return "pptx"
    if ext == ".ppt":
        return "ppt"

    ct = (content_type or "").split(";", 1)[0].strip().lower()
    if ct == "text/plain":
        return "text"
    if ct == "text/markdown":
        return "markdown"
    if ct == "application/pdf":
        return "pdf"
    if ct == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        return "docx"
    if ct == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return "pptx"
    if ct in {"application/msword"}:
        return "doc"
    if ct in {"application/vnd.ms-powerpoint"}:
        return "ppt"
    return ""


def decode_utf8(data: bytes) -> str:
    if data.startswith(b"\xef\xbb\xbf"):
        data = data[3:]
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError as exc:
        raise ValueError("text document is not valid UTF-8") from exc


def parse_docx(data: bytes) -> str:
    from zipfile import ZipFile

    text_parts: list[str] = []
    with ZipFile(io.BytesIO(data)) as zf:
        for name in ("word/document.xml", "word/footnotes.xml", "word/endnotes.xml"):
            if name not in zf.namelist():
                continue
            xml_data = zf.read(name).decode("utf-8", errors="ignore")
            text = extract_text_from_openxml(xml_data)
            if text.strip():
                text_parts.append(text)
    return normalize_text("\n\n".join(text_parts))


def extract_text_from_openxml(xml_text: str) -> str:
    # Rough extraction for docx text nodes.
    text = re.sub(r"<w:tab[^>]*/>", "\t", xml_text)
    text = re.sub(r"<w:br[^>]*/>", "\n", text)
    text = re.sub(r"</w:p>", "\n\n", text)
    text = re.sub(r"<[^>]+>", "", text)
    return text


def parse_pdf(data: bytes) -> tuple[str, list[bytes]]:
    reader = PdfReader(io.BytesIO(data))
    pages: list[str] = []
    images: list[bytes] = []
    for page in reader.pages:
        page_text = (page.extract_text() or "").strip()
        if page_text:
            pages.append(page_text)
        try:
            for image in page.images:
                if image.data:
                    images.append(bytes(image.data))
        except Exception:
            # Keep parsing even if image extraction fails for some pages.
            continue
    return normalize_text("\n\n".join(pages)), images


def parse_pptx(data: bytes) -> tuple[str, list[bytes]]:
    prs = Presentation(io.BytesIO(data))
    texts: list[str] = []
    images: list[bytes] = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        slide_text_parts: list[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                raw_text = (shape.text_frame.text or "").strip()
                if raw_text:
                    slide_text_parts.append(raw_text)
            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    if blob:
                        images.append(blob)
                except Exception:
                    pass
        if slide_text_parts:
            texts.append(f"[Slide {slide_index}]\n" + "\n".join(slide_text_parts))

    return normalize_text("\n\n".join(texts)), images


def build_multimodal_document(title: str, file_type: str, text: str, images: list[bytes]) -> ParsedDocument:
    clean_text = normalize_text(text)
    image_count = len(images)
    heavy = should_use_vision_description(clean_text, image_count)
    used_vision = False
    merged_text = clean_text

    if heavy and image_count > 0:
        vision_desc = describe_images(images[:MAX_DESCRIBE_IMAGES])
        if vision_desc.strip():
            used_vision = True
            if merged_text:
                merged_text = f"{merged_text}\n\n[Image Descriptions]\n{vision_desc}"
            else:
                merged_text = f"[Image Descriptions]\n{vision_desc}"

    return ParsedDocument(
        title=title,
        source_type="TEXT",
        content=ensure_not_empty(merged_text),
        file_type=file_type,
        image_count=image_count,
        used_vision_description=used_vision,
    )


def should_use_vision_description(text: str, image_count: int) -> bool:
    if image_count <= 0:
        return False
    if image_count >= IMAGE_HEAVY_MIN_IMAGES:
        return True
    text_chars = len(text.strip())
    if text_chars == 0:
        return True
    return (text_chars / max(image_count, 1)) < IMAGE_HEAVY_TEXT_CHARS_PER_IMAGE


def describe_images(images: Iterable[bytes]) -> str:
    if not (VISION_BASE_URL and VISION_API_KEY and VISION_MODEL):
        return ""

    descriptions: list[str] = []
    for idx, image_bytes in enumerate(images, start=1):
        desc = describe_single_image(image_bytes, idx)
        if desc:
            descriptions.append(desc)
    return normalize_text("\n\n".join(descriptions))


def describe_single_image(image_bytes: bytes, index: int) -> str:
    data_url = "data:image/png;base64," + base64.b64encode(image_bytes).decode("ascii")
    url = VISION_BASE_URL.rstrip("/") + "/chat/completions"
    payload = {
        "model": VISION_MODEL,
        "messages": [
            {
                "role": "system",
                "content": "You are an OCR and slide-analysis assistant. Describe image content faithfully in concise Chinese.",
            },
            {
                "role": "user",
                "content": [
                    {"type": "text", "text": "请提取这张图里的文字，并简要描述图表/版面重点。输出中文，控制在 150 字以内。"},
                    {"type": "image_url", "image_url": {"url": data_url}},
                ],
            },
        ],
        "temperature": 0.2,
    }
    headers = {
        "Authorization": f"Bearer {VISION_API_KEY}",
        "Content-Type": "application/json",
    }
    try:
        resp = requests.post(url, headers=headers, data=json.dumps(payload), timeout=VISION_TIMEOUT_SECONDS)
        resp.raise_for_status()
        body = resp.json()
        text = body["choices"][0]["message"]["content"]
        return f"[Image {index}] {sanitize_text(str(text))}".strip()
    except Exception:
        return ""


def normalize_text(content: str) -> str:
    content = content.replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.rstrip(" \t") for line in sanitize_text(content).split("\n")]
    cleaned: list[str] = []
    blank = 0
    for line in lines:
        if not line.strip():
            blank += 1
            if blank <= 1:
                cleaned.append("")
            continue
        blank = 0
        cleaned.append(line)
    return "\n".join(cleaned).strip()


def sanitize_text(value: str) -> str:
    out = []
    for ch in value:
        code = ord(ch)
        if code == 0:
            continue
        if ch in ("\t", "\n"):
            out.append(ch)
            continue
        if code < 0x20:
            continue
        out.append(ch)
    return "".join(out)


def ensure_not_empty(text: str) -> str:
    text = (text or "").strip()
    if not text:
        raise ValueError("no extractable text found")
    return text
